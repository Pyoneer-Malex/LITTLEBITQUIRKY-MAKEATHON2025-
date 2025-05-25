import os
import streamlit as st
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation
import random
import re
from collections import Counter
import tempfile
import graphviz
import PyPDF2
import speech_recognition as sr

import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer, WordNetLemmatizer
from transformers import pipeline

# Download NLTK data
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')

# Paths
SCRIPT_DIR = os.path.dirname(__file__)
EXTRACTED_TXT = os.path.join(SCRIPT_DIR, "Extracted.txt")

def write_extracted_txt(text: str):
    """Overwrite Extracted.txt with text."""
    with open(EXTRACTED_TXT, "w", encoding="utf-8") as f:
        f.write(text)

def Cleaner(text):
    filler = {"ok","uh","um","like","yeah","well","hmm","right","actually","just","really","basically"}
    tokens = word_tokenize(text)
    stops = set(stopwords.words("english"))
    filtered = [w.lower() for w in tokens
                if w.isalpha() and w.lower() not in stops and w.lower() not in filler]
    stems = [PorterStemmer().stem(w) for w in filtered]
    return " ".join(WordNetLemmatizer().lemmatize(w) for w in stems)

@st.cache_data
def create_summary(text, max_len=300, min_len=50):
    snippet = text[:3000] if len(text)>3000 else text
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=-1)
    out = summarizer(snippet, max_length=max_len, min_length=min_len, do_sample=False)
    return out[0]["summary_text"].strip()

def generate_yes_no_quiz(cleaned, n=5):
    sents = [s for s in sent_tokenize(cleaned) if len(s.split())>4]
    qs = []
    for s in sents[:n]:
        neg = any(neg in s.lower() for neg in [" not "," no "," never ","n't"])
        qs.append((f"Is it true that: “{s}”", "No" if neg else "Yes"))
    return qs

# --- Streamlit App ---

st.title("🤖 Ο προσωπικός study-buddy σου!")

df = pd.DataFrame(columns=["Αρχείο","Εξαγόμενο Κείμενο"])
uploaded = st.file_uploader("📂 Ανέβασε αρχείο:", type=["png","jpg","jpeg","mp3","pptx","pdf","txt"])

# Session state
st.session_state.setdefault("summary_done", False)
st.session_state.setdefault("summary_text", "")
st.session_state.setdefault("cleaned_summary", "")
st.session_state.setdefault("flashcards", [])
st.session_state.setdefault("fc_idx", 0)
st.session_state.setdefault("show_fc_answer", False)
st.session_state.setdefault("quiz_questions", [])
st.session_state.setdefault("quiz_answers", [])
st.session_state.setdefault("quiz_idx", 0)

extracted = ""
audio_msg = ""
error = False

if uploaded:
    name = uploaded.name.lower()
    try:
        # Image → OCR
        if name.endswith(("png","jpg","jpeg")):
            img = Image.open(uploaded)
            extracted = pytesseract.image_to_string(img)

        # MP3 → speech_recognition
        elif name.endswith("mp3"):
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
            tmp.write(uploaded.read()); tmp.close()
            r = sr.Recognizer()
            with sr.AudioFile(tmp.name) as src:
                audio = r.record(src)
            extracted = r.recognize_google(audio)
            audio_msg = "🎧 MP3 -> κείμενο OK"

        # PPTX → python-pptx
        elif name.endswith("pptx"):
            prs = Presentation(uploaded)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            extracted = "\n".join(lines)
            write_extracted_txt(extracted)

        # PDF → PyPDF2, then write file
        elif name.endswith("pdf"):
            reader = PyPDF2.PdfReader(uploaded)
            pages = [p.extract_text() or "" for p in reader.pages]
            extracted = "\n".join(pages)
            write_extracted_txt(extracted)

        # TXT → read directly, then overwrite Extracted.txt
        elif name.endswith("txt"):
            extracted = uploaded.read().decode("utf-8")
            write_extracted_txt(extracted)

        else:
            st.error("❌ Μη υποστηριζόμενος τύπος.")
            error = True

    except Exception as e:
        st.error(f"❌ Σφάλμα κατά την επεξεργασία: {e}")
        error = True

    if not error:
        df.loc[len(df)] = [uploaded.name, extracted]

        with st.expander("📌 Προεπισκόπηση Κειμένου"):
            st.text_area("Εξαγόμενο Κείμενο", extracted, height=200)
            if audio_msg:
                st.info(audio_msg)

        st.download_button("⬇️ Κατέβασε το κείμενο", extracted, file_name="extracted_text.txt")

        # --- Summarization on the in-memory extracted text ---
        st.subheader("📄 Περίληψη του κειμένου")
        if st.button("📑 Δημιούργησε περίληψη"):
            summary = create_summary(extracted)
            st.session_state.summary_text = summary
            st.session_state.cleaned_summary = Cleaner(extracted)
            st.session_state.summary_done = True

            # prepare quiz
            qas = generate_yes_no_quiz(st.session_state.summary_text, n=5)
            st.session_state.quiz_questions = [q for q,a in qas]
            st.session_state.quiz_answers = [a for q,a in qas]
            st.session_state.quiz_idx = 0

        if st.session_state.summary_done:
            # Display summary as bullets
            st.markdown("**Περίληψη:**")
            for sent in [s.strip() for s in st.session_state.summary_text.split(".") if s.strip()]:
                st.write("• " + sent + ".")

            # --- Flashcards ---
            st.subheader("🃏 Flashcards")
            if not st.session_state.flashcards:
                sents = [s.strip() for s in re.split(r'[.!?]', extracted) if s.strip()]
                cards = []
                for s in random.sample(sents, min(len(sents), 5)):
                    words = [w for w in s.split() if w.isalpha() and len(w) > 3 and w.lower() not in stopwords.words("english")]
                    if not words:
                        continue
                    kw = random.choice(words)
                    pattern = re.compile(re.escape(kw), re.IGNORECASE)
                    masked = pattern.sub("____", s, count=1)
                    cards.append((masked, s, kw))
                st.session_state.flashcards = cards

            if st.session_state.flashcards:
                total_cards = len(st.session_state.flashcards)
                idx = st.session_state.fc_idx
                masked, original, kw = st.session_state.flashcards[idx]

                st.markdown(f"**Κάρτα {idx + 1} από {total_cards}**")
                st.write(f"**Ερώτηση:** {masked}")

                if st.button("🔍 Εμφάνιση απάντησης"):
                    st.session_state.show_fc_answer = True

                if st.session_state.show_fc_answer:
                    pattern = re.compile(re.escape(kw), re.IGNORECASE)
                    bolded = pattern.sub(f"**{kw}**", original, count=1)
                    st.markdown(f"**Απάντηση:** {bolded}")

                col1, col2 = st.columns(2)
                with col1:
                    if st.button("⬅️ Προηγούμενη κάρτα"):
                        st.session_state.show_fc_answer = False
                        st.session_state.fc_idx = (idx - 1) % total_cards
                with col2:
                    if st.button("➡️ Επόμενη κάρτα"):
                        st.session_state.show_fc_answer = False
                        st.session_state.fc_idx = (idx + 1) % total_cards
                        

            # --- Q&A on summary ---
            st.subheader("❓ Ερώτηση πάνω στην περίληψη")
            query = st.text_input("Γράψε λέξη/φράση για Q&A:")
            if query:
                sents = [s for s in sent_tokenize(st.session_state.summary_text) if query.lower() in s.lower()]
                if sents:
                    st.write(sents[0])
                else:
                    st.write("Δεν βρέθηκε πρόταση στην περίληψη με αυτό το στοιχείο.")

            # --- Quiz Yes/No ---
            st.subheader("❓ Quiz Ναι/Όχι")
            idxq = st.session_state.quiz_idx
            if idxq < len(st.session_state.quiz_questions):
                q = st.session_state.quiz_questions[idxq]
                resp = st.radio(q, ("Yes","No"), key=f"quiz_{idxq}")
                if st.button("Υποβολή απάντησης"):
                    correct = st.session_state.quiz_answers[idxq]
                    if resp == correct:
                        st.success("✅ Σωστό")
                    else:
                        st.error(f"❌ Λάθος, σωστό: {correct}")
                    st.session_state.quiz_idx += 1
            else:
                st.write("🏁 Το quiz ολοκληρώθηκε!")

else:
    st.info("Ανέβασε κάποιο αρχείο για να ξεκινήσεις!")

with st.expander("📊 Ιστορικό Αρχείων"):
    st.dataframe(df)